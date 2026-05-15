"""
File content extractor for HTML Presentation Skill.
Reads various file types and outputs structured text for slide generation.

Supported formats: .md, .txt, .pdf, .docx, .pptx, .png, .jpg, .jpeg, .gif, .webp

Usage:
    python read_file.py <file_path> [--output <output_path>]

Output: JSON with extracted content structure.
"""

import sys
import os
import json
import base64
import re
from pathlib import Path


def read_markdown(file_path: str) -> dict:
    """Parse Markdown into sections by headings."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    sections = []
    current_section = {'level': 0, 'title': 'Introduction', 'content': ''}

    for line in content.split('\n'):
        heading_match = re.match(r'^(#{1,6})\s+(.+)', line)
        if heading_match:
            if current_section['content'].strip():
                sections.append(current_section)
            level = len(heading_match.group(1))
            current_section = {
                'level': level,
                'title': heading_match.group(2).strip(),
                'content': ''
            }
        else:
            current_section['content'] += line + '\n'

    if current_section['content'].strip():
        sections.append(current_section)

    # Extract code blocks
    code_blocks = re.findall(r'```(\w*)\n(.*?)```', content, re.DOTALL)

    return {
        'type': 'markdown',
        'file': file_path,
        'title': sections[0]['title'] if sections else Path(file_path).stem,
        'sections': sections,
        'code_blocks': [{'language': lang or 'text', 'code': code.strip()} for lang, code in code_blocks],
        'raw_content': content
    }


def read_text(file_path: str) -> dict:
    """Read plain text file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]

    return {
        'type': 'text',
        'file': file_path,
        'title': Path(file_path).stem,
        'sections': [{'level': 0, 'title': '', 'content': p} for p in paragraphs],
        'raw_content': content
    }


def read_pdf(file_path: str) -> dict:
    """Extract text from PDF using pymupdf (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return {
            'type': 'pdf',
            'file': file_path,
            'error': 'pymupdf not installed. Run: pip install pymupdf',
            'install_cmd': 'pip install pymupdf'
        }

    doc = fitz.open(file_path)
    pages = []
    full_text = ''

    for i, page in enumerate(doc):
        text = page.get_text()
        pages.append({
            'page_number': i + 1,
            'content': text.strip()
        })
        full_text += text + '\n\n'

    # Try to extract sections from headings (larger font = heading)
    sections = []
    for page in pages:
        if page['content']:
            sections.append({
                'level': 1,
                'title': f"Page {page['page_number']}",
                'content': page['content']
            })

    doc.close()

    return {
        'type': 'pdf',
        'file': file_path,
        'title': Path(file_path).stem,
        'page_count': len(pages),
        'sections': sections,
        'raw_content': full_text.strip()
    }


def read_docx(file_path: str) -> dict:
    """Extract text and structure from Word document."""
    try:
        from docx import Document
    except ImportError:
        return {
            'type': 'docx',
            'file': file_path,
            'error': 'python-docx not installed. Run: pip install python-docx',
            'install_cmd': 'pip install python-docx'
        }

    doc = Document(file_path)
    sections = []
    current_section = {'level': 0, 'title': 'Introduction', 'content': ''}

    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            if current_section['content'].strip():
                sections.append(current_section)
            try:
                level = int(para.style.name.replace('Heading ', '').replace('Heading', '1'))
            except ValueError:
                level = 1
            current_section = {
                'level': level,
                'title': para.text.strip(),
                'content': ''
            }
        else:
            current_section['content'] += para.text + '\n'

    if current_section['content'].strip():
        sections.append(current_section)

    # Extract images info
    image_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_count += 1

    return {
        'type': 'docx',
        'file': file_path,
        'title': sections[0]['title'] if sections else Path(file_path).stem,
        'sections': sections,
        'image_count': image_count,
        'raw_content': '\n'.join(p.text for p in doc.paragraphs)
    }


def read_pptx(file_path: str) -> dict:
    """Extract text and structure from PowerPoint."""
    try:
        from pptx import Presentation
    except ImportError:
        return {
            'type': 'pptx',
            'file': file_path,
            'error': 'python-pptx not installed. Run: pip install python-pptx',
            'install_cmd': 'pip install python-pptx'
        }

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': i + 1,
            'title': '',
            'content': '',
            'notes': ''
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data['title'] = shape.text_frame.text
                else:
                    slide_data['content'] += shape.text_frame.text + '\n'

            if hasattr(shape, 'image'):
                slide_data['content'] += f'[Image: {shape.image.content_type}]\n'

        if slide.has_notes_slide:
            slide_data['notes'] = slide.notes_slide.notes_text_frame.text

        slides.append(slide_data)

    return {
        'type': 'pptx',
        'file': file_path,
        'title': slides[0]['title'] if slides else Path(file_path).stem,
        'slide_count': len(slides),
        'slides': slides,
        'raw_content': '\n\n'.join(
            f"## {s['title']}\n{s['content']}" for s in slides
        )
    }


def read_image(file_path: str) -> dict:
    """Read image and return base64-encoded data for embedding."""
    ext = Path(file_path).suffix.lower()
    mime_map = {
        '.png': 'image/png',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.gif': 'image/gif',
        '.webp': 'image/webp',
        '.svg': 'image/svg+xml'
    }

    mime_type = mime_map.get(ext, 'image/png')

    with open(file_path, 'rb') as f:
        data = base64.b64encode(f.read()).decode('utf-8')

    file_size = os.path.getsize(file_path)

    return {
        'type': 'image',
        'file': file_path,
        'title': Path(file_path).stem,
        'mime_type': mime_type,
        'base64_data': data,
        'data_uri': f'data:{mime_type};base64,{data}',
        'file_size_kb': round(file_size / 1024, 1)
    }


READERS = {
    '.md': read_markdown,
    '.markdown': read_markdown,
    '.txt': read_text,
    '.text': read_text,
    '.pdf': read_pdf,
    '.docx': read_docx,
    '.doc': read_docx,   # Note: legacy .doc may not parse correctly; .docx recommended
    '.pptx': read_pptx,
    '.ppt': read_pptx,
    '.png': read_image,
    '.jpg': read_image,
    '.jpeg': read_image,
    '.gif': read_image,
    '.webp': read_image,
    '.svg': read_image,
}


def read_file(file_path: str) -> dict:
    """Auto-detect file type and extract content."""
    path = Path(file_path)

    if not path.exists():
        return {'error': f'File not found: {file_path}'}

    ext = path.suffix.lower()
    reader = READERS.get(ext)

    if not reader:
        return {
            'error': f'Unsupported file type: {ext}',
            'supported': list(READERS.keys())
        }

    try:
        return reader(file_path)
    except Exception as e:
        return {
            'error': f'Failed to read {ext} file: {str(e)}',
            'file': file_path,
            'type': ext.lstrip('.')
        }


def main():
    if len(sys.argv) < 2:
        print("Usage: python read_file.py <file_path> [--output <output_path>]")
        sys.exit(1)

    file_path = sys.argv[1]
    output_path = None

    if '--output' in sys.argv:
        idx = sys.argv.index('--output')
        if idx + 1 < len(sys.argv):
            output_path = sys.argv[idx + 1]

    result = read_file(file_path)

    # For images, don't print base64 to stdout (too large)
    if result.get('type') == 'image' and not output_path:
        display = {k: v for k, v in result.items() if k not in ('base64_data', 'data_uri')}
        display['note'] = 'base64 data omitted from stdout. Use --output to save full JSON.'
        print(json.dumps(display, ensure_ascii=False, indent=2))
    else:
        output = json.dumps(result, ensure_ascii=False, indent=2)
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(output)
            print(f"Saved to {output_path}")
        else:
            print(output)


if __name__ == '__main__':
    main()
